"""
Flask API for PowerPoint Proposal Generator
Supports both Marketing and Software proposals
Vercel-ready with CORS enabled
"""

from flask import Flask, request, jsonify, send_file
from flask_cors import CORS
from pptx import Presentation
from pptx.enum.text import PP_ALIGN
from pptx.util import Pt, Inches
import openai
import re
from datetime import datetime
import json
import smtplib
from email.message import EmailMessage
from io import BytesIO
import stripe
import requests
import os
from dotenv import load_dotenv

load_dotenv()

# ====================
# FLASK APP SETUP
# ====================

app = Flask(__name__)
CORS(app)  # Enable CORS for all routes

# ====================
# CONFIGURATION
# ====================

OPENAI_API_KEY   = os.environ.get("OPENAI_API_KEY")
STRIPE_API_KEY   = os.environ.get("STRIPE_API_KEY")
SENDER_PASSWORD  = os.environ.get("SENDER_PASSWORD")
EXCHANGE_API_KEY = os.environ.get("EXCHANGE_API_KEY")
CALENDLY_TOKEN = os.environ.get("CALENDLY_TOKEN")

client = openai.OpenAI(api_key=OPENAI_API_KEY)
stripe.api_key = STRIPE_API_KEY

# Template configuration - for Vercel deployment
# Files should be in same directory as app.py
TEMPLATES = {
    "marketing": os.path.join(os.path.dirname(__file__), "input.pptx"),
    "software": os.path.join(os.path.dirname(__file__), "input_2.pptx")
}

# Local testing (commented out for deployment)
# TEMPLATES = {
#     "marketing": "input.pptx",
#     "software": "input_2.pptx"
# }

SENDER_EMAIL = "hello@digitaldreamworksagency.com"
REDIRECT_URL = "https://digitaldreamworksagency.com/thank-you"
EXCHANGE_URL = f"https://v6.exchangerate-api.com/v6/{EXCHANGE_API_KEY}/latest/USD"

EVENT_TYPE_URI = "https://api.calendly.com/event_types/6ad037bc-eb04-4132-bc22-66683f8f53d6"

ZERO_DECIMAL_CURRENCIES = ['JPY', 'KRW', 'VND', 'CLP', 'BIF', 'DJF', 'GNF', 'ISK', 'PYG', 'RWF', 'UGX', 'XAF', 'XOF', 'XPF']

# ====================
# CALENDLY FUNCTIONS
# ====================

def generate_calendly_link(client_name, client_email, service_selected, proposal_id=None):
    """Create a Calendly scheduling link with prefilled info"""
    print("📅 Creating personalized Calendly link...")

    url = "https://api.calendly.com/scheduling_links"
    headers = {
        "Authorization": f"Bearer {CALENDLY_TOKEN}",
        "Content-Type": "application/json"
    }

    payload = {
        "max_event_count": 1,
        "owner": EVENT_TYPE_URI,
        "owner_type": "EventType",
        "invitee": {
            "email": client_email,
            "name": client_name
        },
        "metadata": {
            "service": service_selected,
            "proposal_id": proposal_id or "n/a"
        }
    }

    try:
        response = requests.post(url, headers=headers, json=payload)
        response.raise_for_status()
        booking_url = response.json()["resource"]["booking_url"]
        print("✓ Calendly link created:", booking_url)
        return booking_url
    except Exception as e:
        print(f"❌ Calendly error: {e}")
        return "https://calendly.com/digi-dreamworks/onboarding-call"

# ====================
# EMAIL FUNCTIONS
# ====================

def send_email_inmemory(recipient_email, pptx_presentation, project_type):
    """Send an in-memory PPTX as email attachment"""
    try:
        pptx_buffer = BytesIO()
        pptx_presentation.save(pptx_buffer)
        pptx_buffer.seek(0)

        msg = EmailMessage()
        msg['From'] = SENDER_EMAIL
        msg['To'] = recipient_email
        msg['Subject'] = f"Your {project_type.title()} Proposal from Digital Dream Works - {datetime.now().strftime('%B %d, %Y')}"
        msg.set_content(f"Hello,\n\nPlease find attached your personalized {project_type} proposal.\n\nBest regards,\nDigital Dream Works")

        msg.add_attachment(
            pptx_buffer.read(),
            maintype='application',
            subtype='vnd.openxmlformats-officedocument.presentationml.presentation',
            filename=f"{project_type.title()}_Proposal.pptx"
        )

        with smtplib.SMTP_SSL('smtp.gmail.com', 465) as smtp:
            smtp.login(SENDER_EMAIL, SENDER_PASSWORD)
            smtp.send_message(msg)

        print(f"📧 Email sent to {recipient_email} successfully!")
        return True
    except Exception as e:
        print(f"❌ Email error: {e}")
        return False

# ====================
# STRIPE FUNCTIONS
# ====================

def convert_to_usd(amount, currency_code):
    """Convert any currency amount into USD"""
    try:
        print(f"🌍 Converting {amount} {currency_code} to USD...")
        response = requests.get(EXCHANGE_URL)
        data = response.json()

        if "conversion_rates" not in data:
            raise Exception("Bad API response")

        rates = data["conversion_rates"]
        if currency_code.upper() not in rates:
            raise Exception(f"Currency not supported: {currency_code}")

        rate = rates[currency_code.upper()]
        usd_amount = amount / rate
        print(f"   → Converted: ${usd_amount:.2f} USD")
        return usd_amount

    except Exception as e:
        print(f"❌ Currency conversion failed: {e}")
        print("⚠️ Defaulting to USD without conversion.")
        return amount

def create_stripe_payment_link(amount_numeric, currency_code, client_name, service_selected):
    """Create a Stripe Payment Link after converting to USD"""
    print(f"\n💳 Creating Stripe payment link...")
    print(f"   Original Amount: {amount_numeric} {currency_code}")

    try:
        amount_usd = convert_to_usd(amount_numeric, currency_code)
        stripe_amount = int(amount_usd * 100)
        print(f"   Stripe Charge Amount: {stripe_amount} cents (${amount_usd:.2f})")

        product = stripe.Product.create(
            name=f"{service_selected} - {client_name}",
        )

        price = stripe.Price.create(
            unit_amount=stripe_amount,
            currency="usd",
            product=product.id,
        )

        payment_link = stripe.PaymentLink.create(
            line_items=[{"price": price.id, "quantity": 1}],
            after_completion={
                "type": "redirect",
                "redirect": {"url": REDIRECT_URL}
            },
            allow_promotion_codes=False,
        )

        print(f"✓ Payment link created: {payment_link.url}")
        return payment_link.url

    except Exception as e:
        print(f"❌ Error creating payment link: {str(e)}")
        return "https://digitaldreamworksagency.com/payment"

# ====================
# PPT HELPER FUNCTIONS WITH CLICKABLE LINKS
# ====================

from pptx.oxml.xmlchemy import OxmlElement
from pptx.dml.color import RGBColor

def add_hyperlink_to_run(run, url):
    """
    Properly add a hyperlink to a run in PowerPoint.
    This creates an actual clickable link, not just styled text.
    """
    try:
        # Create the hyperlink relationship
        part = run.part
        rId = part.relate_to(
            url,
            "http://schemas.openxmlformats.org/officeDocument/2006/relationships/hyperlink",
            is_external=True
        )
        
        # Get the run's XML element
        r_element = run._r
        
        # Create the hlinkClick element (hyperlink click)
        hlinkClick = OxmlElement('a:hlinkClick')
        hlinkClick.set(
            '{http://schemas.openxmlformats.org/officeDocument/2006/relationships}id',
            rId
        )
        
        # Find or create rPr (run properties)
        rPr = r_element.find('.//a:rPr', namespaces={
            'a': 'http://schemas.openxmlformats.org/drawingml/2006/main'
        })
        
        if rPr is None:
            # Create rPr if it doesn't exist
            rPr = OxmlElement('a:rPr')
            r_element.insert(0, rPr)
        
        # Add the hyperlink to rPr
        rPr.append(hlinkClick)
        
        # Style the text as a link (blue, underlined)
        run.font.color.rgb = RGBColor(0, 0, 255)
        run.font.underline = True
        
        return True
        
    except Exception as e:
        print(f"⚠️ Hyperlink creation failed: {e}")
        # Fallback: just style it as a link
        try:
            run.font.color.rgb = RGBColor(0, 0, 255)
            run.font.underline = True
        except:
            pass
        return False

def find_and_replace_in_text_frame(text_frame, replacements, clickable_links=None):
    """Replace placeholders while preserving formatting and making links clickable"""
    if clickable_links is None:
        clickable_links = {}
    
    for paragraph in text_frame.paragraphs:
        for run in paragraph.runs:
            original_text = run.text
            
            for placeholder, replacement in replacements.items():
                if placeholder in original_text:
                    # Store original formatting
                    font_name = run.font.name
                    font_size = run.font.size
                    font_bold = run.font.bold
                    font_italic = run.font.italic
                    font_color = run.font.color.rgb if run.font.color.type == 1 else None
                    
                    # Replace the text
                    run.text = original_text.replace(placeholder, replacement)
                    
                    # Reapply original formatting
                    if font_name:
                        run.font.name = font_name
                    if font_size:
                        run.font.size = font_size
                    if font_bold is not None:
                        run.font.bold = font_bold
                    if font_italic is not None:
                        run.font.italic = font_italic
                    if font_color:
                        run.font.color.rgb = font_color
                    
                    # Make it clickable if it's a link
                    if placeholder in clickable_links:
                        print(f"🔗 Adding hyperlink to: {placeholder}")
                        add_hyperlink_to_run(run, clickable_links[placeholder])

def find_and_replace_in_shape(shape, replacements, clickable_links=None):
    """Handle different shape types"""
    if hasattr(shape, "text_frame"):
        find_and_replace_in_text_frame(shape.text_frame, replacements, clickable_links)
    elif hasattr(shape, "table"):
        for row in shape.table.rows:
            for cell in row.cells:
                find_and_replace_in_text_frame(cell.text_frame, replacements, clickable_links)
    elif hasattr(shape, "shapes"):
        for sub_shape in shape.shapes:
            find_and_replace_in_shape(sub_shape, replacements, clickable_links)

def replace_placeholders_in_ppt(input_file, replacements, clickable_links=None):
    """Main function to replace all placeholders in PowerPoint and return the presentation object"""
    print(f"\n🔄 Processing: {input_file}\n")
    prs = Presentation(input_file)
    
    for slide_num, slide in enumerate(prs.slides, 1):
        for shape in slide.shapes:
            find_and_replace_in_shape(shape, replacements, clickable_links)
    
    print(f"\n✅ Processing complete")
    return prs

# ====================
# AI GENERATION - SHARED
# ====================

def extract_client_info(email, user_description):
    """Extract structured info from user description"""
    print("\n🤖 Extracting client information...")
    
    prompt = f"""Extract key information from this client inquiry:

Email: {email}
Description: {user_description}

Return as JSON:
{{
  "client_name": "Full name if mentioned, otherwise extract from email",
  "client_company": "Company name if mentioned, otherwise infer or use 'Business'",
  "client_industry": "Best guess of industry",
  "service_selected": "Primary service (1-3 words max)"
}}"""

    response = client.chat.completions.create(
        model="gpt-4o-mini",
        messages=[{"role": "user", "content": prompt}],
        temperature=0.3,
        max_tokens=200,
        response_format={"type": "json_object"}
    )
    
    result = json.loads(response.choices[0].message.content)
    print(f"✓ Extracted: {result['client_name']} | {result['client_industry']} | {result['service_selected']}")
    return result

def calculate_pricing(service_selected, client_industry, user_description):
    """Generate pricing with currency detection"""
    print("\n🤖 Generating: Pricing Estimate...")

    prompt = f"""You are a pricing strategist at Digital Dream Works LLC.

Service: {service_selected}
Industry: {client_industry}
Description: {user_description}

Estimate pricing in the client's local currency. Calculate base price, 15% discount, and final price.

Return as JSON:
{{
  "base_price_display": "X",
  "base_price_numeric": X,
  "currency_code": "X",
  "discount_value_display": "X",
  "discount_value_numeric": X,
  "discounted_price_display": "X",
  "discounted_price_numeric": X,
  "reasoning": "Brief explanation"
}}

CRITICAL:
- currency_code must be ISO 4217 (USD, EUR, GBP, INR, PKR, etc.)
- numeric values are clean numbers (no symbols/commas)
- display values include currency symbols

Price ranges:
- Web Development: $3k-$8k
- AI Marketing: $6k-$15k
- Custom Software/SaaS: $10k-$50k
- Enterprise Solutions: $30k-$100k"""

    response = client.chat.completions.create(
        model="gpt-4o-mini",
        messages=[{"role": "user", "content": prompt}],
        temperature=0.6,
        max_tokens=250,
        response_format={"type": "json_object"}
    )

    result = json.loads(response.choices[0].message.content)
    print(f"✓ Pricing: {result['discounted_price_display']} ({result['currency_code']})")
    return result

# ====================
# AI GENERATION - MARKETING SPECIFIC
# ====================

def generate_marketing_content(client_info, user_description):
    """Generate all marketing-specific content"""
    print("\n🎯 Generating Marketing Proposal Content...")
    
    prompt = f"""You are a marketing strategist at Digital Dream Works LLC.

Client: {client_info['client_name']}
Company: {client_info['client_company']}
Industry: {client_info['client_industry']}
Service: {client_info['service_selected']}
Description: {user_description}

Generate marketing proposal content. Return as JSON:
{{
  "brand_positioning": "1 line, max 180 chars - current market position and digital presence",
  "marketing_systems": "1 line, max 180 chars - current marketing gaps and automation opportunities",
  "growth_potential": "1 line, max 180 chars - growth opportunities with numbers/timeframes",
  "phase1": "Max 150 chars - Foundation Setup phase",
  "phase2": "Max 150 chars - Launch & Scale phase",
  "phase3": "Max 150 chars - Automation & Optimization phase",
  "growth_rate": "20-80 percentage number only",
  "roi_value": "2.0-5.0 multiplier with one decimal",
  "efficiency_rate": "50-85 percentage number only"
}}

Be concise, specific, and actionable."""

    response = client.chat.completions.create(
        model="gpt-4o-mini",
        messages=[{"role": "user", "content": prompt}],
        temperature=0.7,
        max_tokens=600,
        response_format={"type": "json_object"}
    )
    
    result = json.loads(response.choices[0].message.content)
    print("✓ Generated marketing content")
    return result

# ====================
# AI GENERATION - SOFTWARE SPECIFIC
# ====================

def generate_software_content(client_info, user_description):
    """Generate all software-specific content"""
    print("\n💻 Generating Software Proposal Content...")
    
    prompt = f"""You are a software architect at Digital Dream Works LLC.

Client: {client_info['client_name']}
Company: {client_info['client_company']}
Industry: {client_info['client_industry']}
Service: {client_info['service_selected']}
Description: {user_description}

Generate a comprehensive software proposal. Return as JSON:
{{
  "proposal_title": "professional title for the solution. Dont give it a name just identify what it is like maybe an HR chatbot or Chat application etc(max 2 words)",
  "business_problem": "2-3 sentences describing the core business problem",
  "pain_points": "Bullet list of 3-4 key pain points (use • for bullets)",
  "solution_overview": "2-3 sentences paragraph explaining the proposed solution(add 1 • before the para)",
  "solution_value": "2-3 sentences paragraph on business value and ROI (add 1 • before the para)",
  "architecture_description": "2-3 bullet points on system architecture approach (use • for bullets)",
  "module_list": "Bullet list of 4-6 core modules/features (use • for bullets)",
  "tech_stack": "Bullet list of technologies to be used (use • for bullets, e.g., • React.js, • Node.js, • PostgreSQL)",
  "phase1": "Max 150 chars - Requirements & Design phase",
  "phase2": "Max 150 chars - Development & Integration phase",
  "phase3": "Max 150 chars - Testing & QA phase",
  "phase4": "Max 150 chars - Deployment & Training phase",
  "deliverables": "Bullet list of 4-6 key deliverables (use • for bullets)",
  "timeline": "Bullets list matching the deliverables Estimated timeline (use • for bullets)",
  "terms": "Bullets list Brief payment terms (e.g., '30% upfront, 40% milestone, 30% completion' (use • for bullets))",
  "next_steps": "Bullet list of 3 immediate next steps (use • for bullets)"
}}

Be technical yet clear. Focus on business value."""

    response = client.chat.completions.create(
        model="gpt-4o-mini",
        messages=[{"role": "user", "content": prompt}],
        temperature=0.7,
        max_tokens=1200,
        response_format={"type": "json_object"}
    )
    
    result = json.loads(response.choices[0].message.content)
    print("✓ Generated software content")
    return result

# ====================
# MAIN GENERATION PIPELINE
# ====================

def generate_proposal(email, user_description, project_type="marketing"):
    """
    Main function to generate proposal
    
    Args:
        email: Client email
        user_description: Project description
        project_type: "marketing" or "software"
    
    Returns:
        tuple: (pptx_presentation_object, metadata_dict)
    """
    
    if project_type not in TEMPLATES:
        raise ValueError(f"Invalid project_type. Must be 'marketing' or 'software'")
    
    print("\n" + "="*70)
    print(f"🚀 {project_type.upper()} PROPOSAL GENERATION STARTED")
    print("="*70)
    
    # Step 1: Extract client info
    client_info = extract_client_info(email, user_description)
    
    # Step 2: Generate content based on project type
    if project_type == "marketing":
        content = generate_marketing_content(client_info, user_description)
    else:  # software
        content = generate_software_content(client_info, user_description)
    
    # Step 3: Generate pricing
    pricing = calculate_pricing(
        client_info["service_selected"],
        client_info["client_industry"],
        user_description
    )
    
    # Step 4: Generate Stripe payment link
    stripe_link = create_stripe_payment_link(
        amount_numeric=pricing["discounted_price_numeric"],
        currency_code=pricing["currency_code"],
        client_name=client_info["client_name"],
        service_selected=client_info["service_selected"]
    )
    
    # Step 5: Generate Calendly link
    calendly_link = generate_calendly_link(
        client_info["client_name"],
        email,
        client_info["service_selected"]
    )
    
    # Step 6: Build replacements dictionary
    replacements = {
        "{{client_name}}": client_info["client_name"],
        "{{client_company}}": client_info["client_company"],
        "{{base_price}}": pricing["base_price_display"],
        "{{discount}}": pricing["discount_value_display"],
        "{{final_price}}": pricing["discounted_price_display"],
        "{{stripe_link}}": stripe_link,
        "{{calendly_link}}": calendly_link,
        "{{today_date}}": datetime.now().strftime("%B %d, %Y"),
    }
    
    # Clickable links mapping
    clickable_links = {
        "{{stripe_link}}": stripe_link,
        "{{calendly_link}}": calendly_link,
    }
    
    # Add project-specific fields
    if project_type == "marketing":
        replacements.update({
            "{{service_selected}}": client_info["service_selected"],
            "{{client_industry}}": client_info["client_industry"],
            "{{brand_positioning_summary}}": content["brand_positioning"],
            "{{marketing_systems_summary}}": content["marketing_systems"],
            "{{growth_potential_summary}}": content["growth_potential"],
            "{{phase1_summary}}": content["phase1"],
            "{{phase2_summary}}": content["phase2"],
            "{{phase3_summary}}": content["phase3"],
            "{{growth_rate}}": content["growth_rate"],
            "{{roi_value}}": content["roi_value"],
            "{{efficiency_rate}}": content["efficiency_rate"],
            "{{base_price}}": pricing["base_price_display"],
            "{{discount_value}}": pricing["discount_value_display"],
            "{{discounted_price}}": pricing["discounted_price_display"],
            "{{stripe_payment_link}}": stripe_link,
        })
        clickable_links["{{stripe_payment_link}}"] = stripe_link
    else:  # software
        replacements.update({
            "{{proposal_title}}": content["proposal_title"],
            "{{business_problem}}": content["business_problem"],
            "{{pain_points}}": content["pain_points"],
            "{{solution_overview}}": content["solution_overview"],
            "{{solution_value}}": content["solution_value"],
            "{{architecture_description}}": content["architecture_description"],
            "{{module_list}}": content["module_list"],
            "{{tech_stack}}": content["tech_stack"],
            "{{phase1}}": content["phase1"],
            "{{phase2}}": content["phase2"],
            "{{phase3}}": content["phase3"],
            "{{phase4}}": content["phase4"],
            "{{deliverables}}": content["deliverables"],
            "{{timeline}}": content["timeline"],
            "{{terms}}": content["terms"],
            "{{next_steps}}": content["next_steps"],
        })
    
    # Convert all values to strings
    for key, val in replacements.items():
        if isinstance(val, list):
            replacements[key] = "\n".join(val)
        else:
            replacements[key] = str(val)
    
    # Step 7: Generate PPT
    input_file = TEMPLATES[project_type]
    prs = replace_placeholders_in_ppt(input_file, replacements, clickable_links)
    
    # Send email
    send_email_inmemory(email, prs, project_type)
    
    print("\n" + "="*70)
    print("✨ PROPOSAL GENERATION COMPLETE!")
    print("="*70)
    
    metadata = {
        "client_name": client_info["client_name"],
        "client_company": client_info["client_company"],
        "service": client_info["service_selected"],
        "stripe_link": stripe_link,
        "calendly_link": calendly_link,
        "pricing": pricing["discounted_price_display"]
    }
    
    return prs, metadata

# ====================
# FLASK API ROUTES
# ====================

@app.route('/health', methods=['GET'])
def health_check():
    """Health check endpoint"""
    return jsonify({"status": "healthy", "message": "Proposal Generator API is running"}), 200

@app.route('/generate-proposal', methods=['POST'])
def api_generate_proposal():
    """
    API endpoint to generate proposal
    
    Expected JSON body:
    {
        "email": "client@example.com",
        "user_description": "Project details...",
        "project_type": "marketing" or "software"
    }
    """
    try:
        # Get request data
        data = request.get_json()
        
        if not data:
            return jsonify({"error": "No JSON data provided"}), 400
        
        email = data.get('email')
        user_description = data.get('user_description')
        project_type = data.get('project_type', 'marketing')
        
        # Validate inputs
        if not email:
            return jsonify({"error": "Email is required"}), 400
        if not user_description:
            return jsonify({"error": "User description is required"}), 400
        if project_type not in ['marketing', 'software']:
            return jsonify({"error": "Project type must be 'marketing' or 'software'"}), 400
        
        # Generate proposal
        prs, metadata = generate_proposal(email, user_description, project_type)
        
        # Save to BytesIO
        pptx_buffer = BytesIO()
        prs.save(pptx_buffer)
        pptx_buffer.seek(0)
        
        # Return file as download
        filename = f"{project_type.title()}_Proposal_{metadata['client_name'].replace(' ', '_')}.pptx"
        
        return send_file(
            pptx_buffer,
            mimetype='application/vnd.openxmlformats-officedocument.presentationml.presentation',
            as_attachment=True,
            download_name=filename
        )
        
    except Exception as e:
        print(f"❌ API Error: {str(e)}")
        return jsonify({"error": str(e)}), 500

@app.route('/generate-proposal-info', methods=['POST'])
def api_generate_proposal_info():
    """
    API endpoint to generate proposal and return metadata + download link
    (Does not return the file directly, useful for async processing)
    """
    try:
        data = request.get_json()
        
        if not data:
            return jsonify({"error": "No JSON data provided"}), 400
        
        email = data.get('email')
        user_description = data.get('user_description')
        project_type = data.get('project_type', 'marketing')
        
        if not email or not user_description:
            return jsonify({"error": "Email and user_description are required"}), 400
        
        if project_type not in ['marketing', 'software']:
            return jsonify({"error": "Project type must be 'marketing' or 'software'"}), 400
        
        # Generate proposal
        prs, metadata = generate_proposal(email, user_description, project_type)
        
        return jsonify({
            "success": True,
            "message": "Proposal generated and sent via email",
            "metadata": metadata
        }), 200
        
    except Exception as e:
        print(f"❌ API Error: {str(e)}")
        return jsonify({"error": str(e)}), 500

# ====================
# VERCEL SERVERLESS HANDLER
# ====================

# For Vercel deployment
def handler(request):
    with app.app_context():
        return app(request)

# ====================
# LOCAL TESTING
# ====================

if __name__ == "__main__":
    # Local development server
    app.run(debug=True, host='0.0.0.0', port=5000)